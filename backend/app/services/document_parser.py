"""
Document Parser Service
Supports: PDF, DOCX, PPTX, XLSX, TXT, HTML, MD, JSON, CSV
"""
import io
import logging
from typing import Dict, Any, Optional
from pathlib import Path

# Document parsing libraries
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import markdown
import chardet
import json
import csv

logger = logging.getLogger(__name__)


class DocumentParserService:
    """Service for parsing various document formats"""

    # Supported MIME types
    SUPPORTED_MIME_TYPES = {
        # PDF
        "application/pdf": "pdf",
        # Microsoft Office
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/msword": "doc",
        # Text formats
        "text/plain": "txt",
        "text/html": "html",
        "text/markdown": "md",
        "text/csv": "csv",
        "application/json": "json",
        # Email
        "message/rfc822": "eml",
    }

    # Maximum file size (50 MB)
    MAX_FILE_SIZE = 50 * 1024 * 1024

    def __init__(self):
        """Initialize document parser"""
        self.logger = logger

    async def parse_document(
        self,
        content: bytes,
        mime_type: str,
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Parse document and extract text content

        Args:
            content: Document content as bytes
            mime_type: MIME type of the document
            filename: Optional filename for extension detection

        Returns:
            Dict with parsed content and metadata

        Raises:
            ValueError: If document type is not supported or parsing fails
        """
        # Check file size
        if len(content) > self.MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds maximum of {self.MAX_FILE_SIZE} bytes")

        # Determine document type
        doc_type = self._get_document_type(mime_type, filename)

        if not doc_type:
            raise ValueError(f"Unsupported document type: {mime_type}")

        self.logger.info(f"Parsing document type: {doc_type}, size: {len(content)} bytes")

        # Parse based on type
        try:
            if doc_type == "pdf":
                return await self._parse_pdf(content)
            elif doc_type == "docx":
                return await self._parse_docx(content)
            elif doc_type == "pptx":
                return await self._parse_pptx(content)
            elif doc_type == "xlsx":
                return await self._parse_xlsx(content)
            elif doc_type == "txt":
                return await self._parse_text(content)
            elif doc_type == "html":
                return await self._parse_html(content)
            elif doc_type == "md":
                return await self._parse_markdown(content)
            elif doc_type == "json":
                return await self._parse_json(content)
            elif doc_type == "csv":
                return await self._parse_csv(content)
            else:
                raise ValueError(f"Parser not implemented for type: {doc_type}")

        except Exception as e:
            self.logger.error(f"Error parsing document: {str(e)}", exc_info=True)
            raise ValueError(f"Failed to parse document: {str(e)}")

    def _get_document_type(self, mime_type: str, filename: Optional[str] = None) -> Optional[str]:
        """Determine document type from MIME type or filename"""

        # Try MIME type first
        if mime_type in self.SUPPORTED_MIME_TYPES:
            return self.SUPPORTED_MIME_TYPES[mime_type]

        # Try filename extension
        if filename:
            ext = Path(filename).suffix.lower().lstrip(".")
            if ext in ["pdf", "docx", "pptx", "xlsx", "txt", "html", "md", "json", "csv"]:
                return ext

        return None

    async def _parse_pdf(self, content: bytes) -> Dict[str, Any]:
        """Parse PDF document"""
        pdf_file = io.BytesIO(content)
        reader = pypdf.PdfReader(pdf_file)

        # Extract text from all pages
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)

        full_text = "\n\n".join(text_parts)

        # Get metadata
        metadata = {}
        if reader.metadata:
            metadata = {
                "title": reader.metadata.get("/Title"),
                "author": reader.metadata.get("/Author"),
                "subject": reader.metadata.get("/Subject"),
                "creator": reader.metadata.get("/Creator"),
                "producer": reader.metadata.get("/Producer"),
                "creation_date": reader.metadata.get("/CreationDate"),
            }

        return {
            "content": full_text,
            "word_count": len(full_text.split()),
            "char_count": len(full_text),
            "page_count": len(reader.pages),
            "metadata": metadata,
        }

    async def _parse_docx(self, content: bytes) -> Dict[str, Any]:
        """Parse DOCX document"""
        docx_file = io.BytesIO(content)
        doc = DocxDocument(docx_file)

        # Extract text from paragraphs
        text_parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        full_text = "\n\n".join(text_parts)

        # Get metadata
        metadata = {
            "title": doc.core_properties.title,
            "author": doc.core_properties.author,
            "subject": doc.core_properties.subject,
            "created": str(doc.core_properties.created) if doc.core_properties.created else None,
            "modified": str(doc.core_properties.modified) if doc.core_properties.modified else None,
        }

        return {
            "content": full_text,
            "word_count": len(full_text.split()),
            "char_count": len(full_text),
            "paragraph_count": len(doc.paragraphs),
            "metadata": metadata,
        }

    async def _parse_pptx(self, content: bytes) -> Dict[str, Any]:
        """Parse PPTX presentation"""
        pptx_file = io.BytesIO(content)
        prs = Presentation(pptx_file)

        # Extract text from all slides
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

        full_text = "\n\n".join(text_parts)

        # Get metadata
        metadata = {
            "title": prs.core_properties.title,
            "author": prs.core_properties.author,
            "subject": prs.core_properties.subject,
            "created": str(prs.core_properties.created) if prs.core_properties.created else None,
        }

        return {
            "content": full_text,
            "word_count": len(full_text.split()),
            "char_count": len(full_text),
            "slide_count": len(prs.slides),
            "metadata": metadata,
        }

    async def _parse_xlsx(self, content: bytes) -> Dict[str, Any]:
        """Parse XLSX spreadsheet"""
        xlsx_file = io.BytesIO(content)
        wb = load_workbook(xlsx_file, data_only=True)

        # Extract text from all sheets
        text_parts = []
        for sheet in wb.worksheets:
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    sheet_text.append(" | ".join(row_text))

            if sheet_text:
                text_parts.append(f"Sheet: {sheet.title}\n" + "\n".join(sheet_text))

        full_text = "\n\n".join(text_parts)

        metadata = {
            "sheet_count": len(wb.worksheets),
            "sheet_names": [sheet.title for sheet in wb.worksheets],
        }

        return {
            "content": full_text,
            "word_count": len(full_text.split()),
            "char_count": len(full_text),
            "metadata": metadata,
        }

    async def _parse_text(self, content: bytes) -> Dict[str, Any]:
        """Parse plain text document"""
        # Detect encoding
        detection = chardet.detect(content)
        encoding = detection["encoding"] or "utf-8"

        # Decode text
        text = content.decode(encoding, errors="ignore")

        return {
            "content": text,
            "word_count": len(text.split()),
            "char_count": len(text),
            "metadata": {"encoding": encoding},
        }

    async def _parse_html(self, content: bytes) -> Dict[str, Any]:
        """Parse HTML document"""
        # Detect encoding
        detection = chardet.detect(content)
        encoding = detection["encoding"] or "utf-8"

        # Parse HTML
        soup = BeautifulSoup(content, "lxml")

        # Remove script and style tags
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text
        text = soup.get_text(separator="\n", strip=True)

        # Get title
        title = soup.title.string if soup.title else None

        return {
            "content": text,
            "word_count": len(text.split()),
            "char_count": len(text),
            "metadata": {
                "title": title,
                "encoding": encoding,
            },
        }

    async def _parse_markdown(self, content: bytes) -> Dict[str, Any]:
        """Parse Markdown document"""
        # Detect encoding
        detection = chardet.detect(content)
        encoding = detection["encoding"] or "utf-8"

        # Decode markdown
        md_text = content.decode(encoding, errors="ignore")

        # Convert to HTML and extract text
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, "lxml")
        text = soup.get_text(separator="\n", strip=True)

        return {
            "content": text,
            "word_count": len(text.split()),
            "char_count": len(text),
            "metadata": {"encoding": encoding},
        }

    async def _parse_json(self, content: bytes) -> Dict[str, Any]:
        """Parse JSON document"""
        # Detect encoding
        detection = chardet.detect(content)
        encoding = detection["encoding"] or "utf-8"

        # Parse JSON
        json_data = json.loads(content.decode(encoding))

        # Convert to readable text
        text = json.dumps(json_data, indent=2)

        return {
            "content": text,
            "word_count": len(text.split()),
            "char_count": len(text),
            "metadata": {"encoding": encoding},
        }

    async def _parse_csv(self, content: bytes) -> Dict[str, Any]:
        """Parse CSV document"""
        # Detect encoding
        detection = chardet.detect(content)
        encoding = detection["encoding"] or "utf-8"

        # Parse CSV
        csv_text = content.decode(encoding, errors="ignore")
        csv_reader = csv.reader(io.StringIO(csv_text))

        # Convert to text
        rows = []
        for row in csv_reader:
            if row:
                rows.append(" | ".join(row))

        text = "\n".join(rows)

        return {
            "content": text,
            "word_count": len(text.split()),
            "char_count": len(text),
            "metadata": {
                "encoding": encoding,
                "row_count": len(rows),
            },
        }

    def chunk_text(
        self,
        text: str,
        chunk_size: int = 1000,
        overlap: int = 200,
    ) -> list[Dict[str, Any]]:
        """
        Split text into overlapping chunks for better vector search

        Args:
            text: Text to chunk
            chunk_size: Maximum characters per chunk
            overlap: Number of characters to overlap between chunks

        Returns:
            List of chunks with metadata
        """
        chunks = []

        # If text is smaller than chunk size, return as single chunk
        if len(text) <= chunk_size:
            return [{
                "index": 0,
                "content": text,
                "char_count": len(text),
            }]

        # Split into chunks with overlap
        start = 0
        index = 0

        while start < len(text):
            end = min(start + chunk_size, len(text))

            # Try to break at sentence or word boundary
            if end < len(text):
                # Look for sentence end
                for delimiter in [".", "!", "?", "\n"]:
                    last_delimiter = text[start:end].rfind(delimiter)
                    if last_delimiter > chunk_size * 0.5:  # At least 50% of chunk size
                        end = start + last_delimiter + 1
                        break
                else:
                    # Look for word boundary
                    last_space = text[start:end].rfind(" ")
                    if last_space > chunk_size * 0.5:
                        end = start + last_space + 1

            chunk_text = text[start:end].strip()

            if chunk_text:
                chunks.append({
                    "index": index,
                    "content": chunk_text,
                    "char_count": len(chunk_text),
                })
                index += 1

            # Move to next chunk with overlap
            start = end - overlap if end < len(text) else end

        return chunks


# Create singleton instance
document_parser_service = DocumentParserService()
